from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

NAVY = RGBColor(0x0F, 0x2A, 0x4A)
TEAL = RGBColor(0x0E, 0x7C, 0x86)
ORANGE = RGBColor(0xE5, 0x6B, 0x1F)
LIGHT_GREY = RGBColor(0xEE, 0xEE, 0xEE)
DARK_GREY = RGBColor(0x33, 0x33, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_text(text_frame, text, *, size=18, bold=False, color=DARK_GREY, align=None):
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    if align is not None:
        paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(slide, left, top, width, height, text, *, size=14, bold=False, color=DARK_GREY, align=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    box.text_frame.word_wrap = True
    box.text_frame.margin_left = Inches(0.05)
    box.text_frame.margin_top = Inches(0.02)
    set_text(box.text_frame, text, size=size, bold=bold, color=color, align=align)
    return box


def add_bullets(slide, left, top, width, height, bullets, *, size=12, color=DARK_GREY):
    box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = box.text_frame
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.05)
    text_frame.margin_top = Inches(0.02)
    text_frame.clear()
    for index, bullet in enumerate(bullets):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        run = paragraph.add_run()
        run.text = f"•  {bullet}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        paragraph.space_after = Pt(1)
    return box


def add_pill(slide, left, top, width, height, label, fill, *, size=11, bold=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_top = Inches(0.04)
    set_text(tf, label, size=size, bold=bold, color=WHITE, align=PP_ALIGN.CENTER)
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=NAVY, width=1.5):
    connector = slide.shapes.add_connector(2, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    return connector


def add_section_header(slide, left, top, width, text):
    box = slide.shapes.add_textbox(left, top, width, Inches(0.3))
    box.text_frame.margin_left = Inches(0.05)
    set_text(box.text_frame, text, size=13, bold=True, color=NAVY)


def main() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.6))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = NAVY
    title_bar.line.fill.background()
    set_text(title_bar.text_frame, "SSL Sustainable Solutions Lab Chatbot — System Overview", size=20, bold=True, color=WHITE)
    title_bar.text_frame.margin_left = Inches(0.3)
    title_bar.text_frame.margin_top = Inches(0.12)

    # ---- Top row: Request flow pipeline ----
    add_section_header(slide, Inches(0.3), Inches(0.75), Inches(12.7), "Request flow — from user question to streamed answer")

    pipeline_steps = [
        ("1. Receive", "Flask /api/chat\nuser_message + history", TEAL),
        ("2. Resolve", "follow-up resolver\n+ generic anchor", TEAL),
        ("3. Route", "detect_local_query_route\nintent classifier", TEAL),
        ("4. Fast path", "section / entity /\ndocument registries", NAVY),
        ("5. Retrieve", "Chroma dense + BM25\nfused via RRF + rerank", NAVY),
        ("6. Plan?", "low-conf: plan_query_with_llm\n(Gemini, budget=0)", ORANGE),
        ("7. Stream", "Gemini main answer\ngenerate_content_stream", ORANGE),
        ("8. Done", "filter_sources_to_cited\n+ attach_trace + SSE", TEAL),
    ]
    pipe_top = Inches(1.1)
    pipe_height = Inches(0.85)
    pipe_x_start = Inches(0.3)
    pipe_total_width = Inches(12.7)
    gap = Inches(0.08)
    step_width = (pipe_total_width - gap * (len(pipeline_steps) - 1)) / len(pipeline_steps)
    centers_y = pipe_top + pipe_height / 2
    for index, (title, detail, color) in enumerate(pipeline_steps):
        left = pipe_x_start + (step_width + gap) * index
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, pipe_top, step_width, pipe_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = color
        shape.line.width = Pt(1.5)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.05)
        tf.margin_top = Inches(0.05)
        tf.margin_right = Inches(0.05)
        tf.clear()
        t_run = tf.paragraphs[0].add_run()
        t_run.text = title
        t_run.font.size = Pt(11)
        t_run.font.bold = True
        t_run.font.color.rgb = color
        d_para = tf.add_paragraph()
        d_run = d_para.add_run()
        d_run.text = detail
        d_run.font.size = Pt(9)
        d_run.font.color.rgb = DARK_GREY
        if index < len(pipeline_steps) - 1:
            add_arrow(slide, left + step_width, centers_y, left + step_width + gap, centers_y)

    # Parallel suggestions callout
    sugg_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(2.05), Inches(12.7), Inches(0.4))
    sugg_box.fill.solid()
    sugg_box.fill.fore_color.rgb = LIGHT_GREY
    sugg_box.line.color.rgb = ORANGE
    sugg_box.line.width = Pt(1.25)
    set_text(
        sugg_box.text_frame,
        "Parallel thread:  generate_suggestions (Gemini, thinking_budget=0)  starts at step 4 and is joined after step 7 — hides behind the answer stream.",
        size=11,
        bold=True,
        color=DARK_GREY,
        align=PP_ALIGN.CENTER,
    )
    sugg_box.text_frame.margin_top = Inches(0.07)

    # ---- Bottom area: three columns ----
    columns_top = Inches(2.65)
    col_width = Inches(4.2)
    col_height = Inches(4.65)
    col_gap = Inches(0.12)
    col_x_start = Inches(0.3)

    # Column 1: Retrieval
    col1_left = col_x_start
    col1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, col1_left, columns_top, col_width, col_height)
    col1.fill.solid()
    col1.fill.fore_color.rgb = WHITE
    col1.line.color.rgb = TEAL
    col1.line.width = Pt(1.75)
    add_text(slide, col1_left + Inches(0.15), columns_top + Inches(0.08), col_width - Inches(0.3), Inches(0.3),
             "Retrieval & Indexing", size=14, bold=True, color=TEAL)
    add_bullets(
        slide,
        col1_left + Inches(0.15),
        columns_top + Inches(0.45),
        col_width - Inches(0.3),
        col_height - Inches(0.5),
        [
            "Corpus: SEED_DOCUMENTS (.txt + .pdf)",
            "Chunking: langchain_text_splitters",
            "Embeddings: sentence-transformers (local)",
            "Vector store: Chroma (./chroma_db)",
            "Lexical: rank_bm25 (in-process BM25)",
            "Fusion: Reciprocal Rank Fusion (dense + BM25)",
            "Rerank: boosts on target_source_paths / titles /",
            "    folders / categories, summary vs detail, exact-",
            "    token matches in section names",
            "Three fast-path registries short-circuit retrieval:",
            "    section_registry, entity_registry,",
            "    document_registry",
        ],
        size=10,
    )

    # Column 2: LLM Layer
    col2_left = col_x_start + (col_width + col_gap)
    col2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, col2_left, columns_top, col_width, col_height)
    col2.fill.solid()
    col2.fill.fore_color.rgb = WHITE
    col2.line.color.rgb = ORANGE
    col2.line.width = Pt(1.75)
    add_text(slide, col2_left + Inches(0.15), columns_top + Inches(0.08), col_width - Inches(0.3), Inches(0.3),
             "LLM Layer — Gemini API", size=14, bold=True, color=ORANGE)
    add_bullets(
        slide,
        col2_left + Inches(0.15),
        columns_top + Inches(0.45),
        col_width - Inches(0.3),
        col_height - Inches(0.5),
        [
            "Provider: Google Gemini API",
            "SDK: google.genai",
            "Model: gemini-3.1-flash-lite",
            "Calls per request (up to 3):",
            "    • plan_query_with_llm — only when retrieval",
            "        is low-confidence; budget=0,  ~3–6 s",
            "    • Main answer — always (unless fast path);",
            "        streamed via generate_content_stream,",
            "        ~5–10 s",
            "    • generate_suggestions — parallel thread,",
            "        budget=0, ~4–8 s (hidden by stream)",
            "Retry: 429 / 503 / UNAVAILABLE / timeouts retried",
            "    up to 8 attempts with API-supplied delays",
        ],
        size=10,
    )

    # Column 3: Stack & Endpoints
    col3_left = col_x_start + 2 * (col_width + col_gap)
    col3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, col3_left, columns_top, col_width, col_height)
    col3.fill.solid()
    col3.fill.fore_color.rgb = WHITE
    col3.line.color.rgb = NAVY
    col3.line.width = Pt(1.75)
    add_text(slide, col3_left + Inches(0.15), columns_top + Inches(0.08), col_width - Inches(0.3), Inches(0.3),
             "Stack & Endpoints", size=14, bold=True, color=NAVY)
    add_bullets(
        slide,
        col3_left + Inches(0.15),
        columns_top + Inches(0.45),
        col_width - Inches(0.3),
        col_height - Inches(0.5),
        [
            "Backend: Python 3.9, Flask, python-dotenv",
            "Frontend: vanilla JS + fetch stream reader,",
            "    HTML / CSS templates",
            "Content filter: better_profanity (pre-LLM)",
            "PDF ingestion: pypdf",
            "Analytics: local /dashboard with trace events",
            "",
            "Endpoints (Flask):",
            "    GET  /                         chat UI",
            "    GET  /dashboard                analytics",
            "    GET  /dashboard/interactions/<id>",
            "    POST /api/chat                 SSE stream:",
            "         meta / delta / suggestions / done",
        ],
        size=10,
    )

    # Footer
    foot = slide.shapes.add_textbox(Inches(0.3), Inches(7.15), Inches(12.7), Inches(0.3))
    set_text(foot.text_frame, "RAG over the SSL document corpus  ·  Flask + Chroma + Gemini  ·  Streaming SSE with parallel suggestion generation",
             size=10, color=NAVY, align=PP_ALIGN.CENTER)

    out_path = Path(__file__).resolve().parent / "SSL_Chatbot_System_Overview.pptx"
    prs.save(out_path)
    print(f"Wrote {out_path}")


if __name__ == "__main__":
    main()
